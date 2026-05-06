import os
import tempfile
import requests
import fitz
from docx import Document
from pptx import Presentation
from PIL import Image


OCR_SEND_LIMIT_BYTES = 950 * 1024  # compress before sending to OCR demo API


def extract_text(uploaded_file, use_online_ocr=True, ocr_api_key="helloworld"):
    ext = os.path.splitext(uploaded_file.name)[1].lower()

    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(uploaded_file.getbuffer())
        path = tmp.name

    try:
        if ext == ".pdf":
            return extract_pdf(path, use_online_ocr, ocr_api_key)

        if ext == ".docx":
            return extract_docx(path)

        if ext == ".pptx":
            return extract_pptx(path)

        if ext in [".png", ".jpg", ".jpeg"]:
            if not use_online_ocr:
                return ""
            compressed_path = compress_image_for_ocr(path)
            return ocr_space_file(compressed_path, ocr_api_key)

        return ""

    finally:
        try:
            os.remove(path)
        except OSError:
            pass


def extract_pdf(path, use_online_ocr=True, ocr_api_key="helloworld"):
    doc = fitz.open(path)
    text_parts = []

    for page in doc:
        txt = page.get_text("text", sort=True).strip()
        if txt:
            text_parts.append(txt)

    if text_parts:
        return "\n\n".join(text_parts).strip()

    if use_online_ocr:
        # OCR.Space boleh membaca PDF scan, tetapi akaun demo/free mungkin ada had saiz.
        return ocr_space_file(path, ocr_api_key, is_pdf=True)

    return ""


def extract_docx(path):
    doc = Document(path)
    parts = []

    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = " ".join(
                cell.text.strip()
                for cell in row.cells
                if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)

    return "\n".join(parts).strip()


def extract_pptx(path):
    prs = Presentation(path)
    parts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())

    return "\n".join(parts).strip()


def compress_image_for_ocr(path):
    """
    Terima imej sehingga 10 MB, kemudian kecilkan kepada fail JPEG bawah ±950 KB
    sebelum dihantar ke OCR.Space demo/free API.
    """
    image = Image.open(path).convert("RGB")

    max_width = 1800
    if image.width > max_width:
        ratio = max_width / image.width
        new_height = int(image.height * ratio)
        image = image.resize((max_width, new_height))

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
    tmp.close()

    quality = 90
    while quality >= 35:
        image.save(tmp.name, format="JPEG", quality=quality, optimize=True)
        if os.path.getsize(tmp.name) <= OCR_SEND_LIMIT_BYTES:
            return tmp.name
        quality -= 10

    # Jika masih besar, resize lagi.
    while os.path.getsize(tmp.name) > OCR_SEND_LIMIT_BYTES and image.width > 900:
        new_width = int(image.width * 0.85)
        new_height = int(image.height * 0.85)
        image = image.resize((new_width, new_height))
        image.save(tmp.name, format="JPEG", quality=45, optimize=True)

    return tmp.name


def ocr_space_file(path, api_key="helloworld", is_pdf=False):
    url = "https://api.ocr.space/parse/image"

    payload = {
        "apikey": api_key or "helloworld",
        "language": "eng",
        "isOverlayRequired": False,
        "OCREngine": 2,
        "scale": True,
        "detectOrientation": True,
    }

    if is_pdf:
        payload["isCreateSearchablePdf"] = False

    with open(path, "rb") as f:
        response = requests.post(
            url,
            files={"filename": f},
            data=payload,
            timeout=60
        )

    if response.status_code != 200:
        raise RuntimeError(f"OCR API gagal. Status code: {response.status_code}")

    result = response.json()

    if result.get("IsErroredOnProcessing"):
        message = result.get("ErrorMessage") or result.get("ErrorDetails") or "OCR gagal."
        raise RuntimeError(str(message))

    parsed = result.get("ParsedResults") or []

    texts = []
    for item in parsed:
        text = item.get("ParsedText", "")
        if text.strip():
            texts.append(text.strip())

    return "\n\n".join(texts).strip()
